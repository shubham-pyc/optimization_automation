from pptx import Presentation
from campaign import ExceptionHander
class TextMinner:
    def __init__(self,file_path):
        self.exception_handler = ExceptionHander()
        self.presentation = None
        if isinstance(file_path,str):
            try:
                self.file = open(file_path)
                self.presentation = Presentation(self.file)
            except Exception as e:
                self.exception_handler.handle_exception(e)
        else:
            raise ValueError("Invalid file for Presentation")
    
    def extract_text(self):
        ret_value = []
        try:
            for slide in self.presentation.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            ret_value.append(run.text)

            print(ret_value)
        except Exception as e:
            self.exception_handler.handle_exception(e)
    
    def mine_text(self,array):
        pass

        
minner = TextMinner("Suppress+-3 day v1.pptx")
minner.extract_text()
